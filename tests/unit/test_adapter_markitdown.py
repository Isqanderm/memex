import pytest
from src.adapters.protocol import Source
from src.adapters.markitdown_adapter import (
    MarkItDownAdapter,
    PPTX_MIME, XLSX_MIME, EPUB_MIME,
)


def test_can_handle_pptx_by_mime():
    assert MarkItDownAdapter().can_handle(Source(path="deck.pptx", mime_type=PPTX_MIME))


def test_can_handle_xlsx_by_mime():
    assert MarkItDownAdapter().can_handle(Source(path="data.xlsx", mime_type=XLSX_MIME))


def test_can_handle_epub_by_mime():
    assert MarkItDownAdapter().can_handle(Source(path="book.epub", mime_type=EPUB_MIME))


def test_can_handle_by_extension_pptx():
    assert MarkItDownAdapter().can_handle(Source(path="deck.pptx", mime_type="application/octet-stream"))


def test_can_handle_by_extension_xlsx():
    assert MarkItDownAdapter().can_handle(Source(path="data.xlsx", mime_type="application/octet-stream"))


def test_can_handle_by_extension_epub():
    assert MarkItDownAdapter().can_handle(Source(path="book.epub", mime_type="application/octet-stream"))


def test_cannot_handle_pdf():
    assert not MarkItDownAdapter().can_handle(Source(path="doc.pdf", mime_type="application/pdf"))


def test_cannot_handle_txt():
    assert not MarkItDownAdapter().can_handle(Source(path="doc.txt", mime_type="text/plain"))


def test_parses_xlsx(tmp_path):
    pytest.importorskip("openpyxl")
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alpha", 1])
    ws.append(["Beta", 2])
    path = tmp_path / "test.xlsx"
    wb.save(str(path))

    doc = MarkItDownAdapter().parse(Source(path=str(path), mime_type=XLSX_MIME))
    assert doc.mime_type == XLSX_MIME
    assert len(doc.sections) >= 1
    combined = " ".join(s.content for s in doc.sections)
    assert "Alpha" in combined or "Sheet1" in combined


def test_parses_pptx(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "My Slide"
    slide.placeholders[1].text = "Slide content here."
    path = tmp_path / "test.pptx"
    prs.save(str(path))

    doc = MarkItDownAdapter().parse(Source(path=str(path), mime_type=PPTX_MIME))
    assert doc.mime_type == PPTX_MIME
    assert len(doc.sections) >= 1
    combined = " ".join(s.content for s in doc.sections)
    assert "My Slide" in combined or "Slide content" in combined


def test_returns_parsed_document_structure(tmp_path):
    pytest.importorskip("openpyxl")
    import openpyxl
    from src.models.parsed import ParsedDocument
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Col"])
    ws.append(["Val"])
    path = tmp_path / "mini.xlsx"
    wb.save(str(path))

    doc = MarkItDownAdapter().parse(Source(path=str(path), mime_type=XLSX_MIME))
    assert isinstance(doc, ParsedDocument)
    assert doc.source == str(path)
    assert all(s.content for s in doc.sections)
